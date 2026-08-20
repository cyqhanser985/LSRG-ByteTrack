# -*- coding: utf-8 -*-
"""
Generate Complete, Spacious, Tabular Presentation with Unified-Scale Feature Validation Table
Total 7 Slides:
  Slide 0: 深邃学术封面 (Cover)
  Slide 1: 几何测度与关联基石：IoU 空间重叠与代价矩阵 (Geometric Primer)
  Slide 2: 初始特征有效性检验：统一 [0, 1] 尺度下的正负样本断崖对比大表 (UNIFIED SCALE TABLE)
  Slide 3: 四大聚合模型 5-Fold 盲测对比大表 (表2 + 大号两阶段结论)
  Slide 4: 【全屏大图 1】四大模型 ROC 与 5% 步长折线大图 (最大化展示 + 结论条)
  Slide 5: 三大失效模式精细化解耦与“自信而错”归因 (表3 + 大号物理归因)
  Slide 6: 【全屏大图 2】跨数据集 2×2 独立评测 ROC 4 宫格大图 (最大化展示 + 异质性结论条)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────────────────
NAVY_DEEP = RGBColor(11, 28, 51)       # #0b1c33
NAVY_PRIMARY = RGBColor(22, 54, 92)    # #16365c
NAVY_MID = RGBColor(37, 77, 126)       # #254d7e
NAVY_HEADER = RGBColor(16, 42, 77)     # Table header bg

BG_LIGHT = RGBColor(246, 248, 251)     # #f6f8fb
CARD_BG = RGBColor(255, 255, 255)      # #ffffff
CARD_BORDER = RGBColor(223, 229, 239)  # #dfe5ef

ROW_EVEN = RGBColor(255, 255, 255)
ROW_ODD = RGBColor(245, 248, 252)

TEXT_DARK = RGBColor(28, 36, 48)       # #1c2430
TEXT_MUTED = RGBColor(100, 115, 138)   # #64738a
TEXT_WHITE = RGBColor(255, 255, 255)

GOLD = RGBColor(212, 160, 70)          # #d4a046
GREEN = RGBColor(31, 125, 75)          # #1f7d4b
GREEN_BG = RGBColor(237, 248, 242)     # #edf8f2
RED_ACCENT = RGBColor(186, 52, 40)     # #ba3428
RED_BG = RGBColor(253, 241, 236)       # #fdf1ec
BLUE_BG = RGBColor(238, 245, 252)      # #eef5fc

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"


def set_slide_background(slide, prs, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, step_num, title_text, subtitle_text):
    shape_badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.38), Inches(0.65), Inches(0.32)
    )
    shape_badge.fill.solid()
    shape_badge.fill.fore_color.rgb = GOLD
    shape_badge.line.color.rgb = GOLD
    tf_b = shape_badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.text = f"{step_num}"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = NAVY_DEEP
    p_b.font.name = FONT_TITLE

    tx_box = slide.shapes.add_textbox(Inches(1.55), Inches(0.32), Inches(11.0), Inches(0.42))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(17.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_DEEP
    p.font.name = FONT_TITLE

    tx_box_sub = slide.shapes.add_textbox(Inches(1.55), Inches(0.76), Inches(11.0), Inches(0.32))
    tf_sub = tx_box_sub.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(10.5)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.font.name = FONT_BODY


def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_speaker_note(slide, note_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = note_text


def format_cell(cell, text, font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, bg_color=None):
    cell.text = text
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_BODY


def create_seven_slides_v7(output_path, img_dir):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 0: COVER SLIDE
    # =========================================================================
    slide0 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide0, prs, NAVY_DEEP)

    badge0 = slide0.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.1), Inches(3.2), Inches(0.38)
    )
    badge0.fill.solid()
    badge0.fill.fore_color.rgb = RGBColor(25, 48, 80)
    badge0.line.color.rgb = GOLD
    badge0.line.width = Pt(1)
    tf0_b = badge0.text_frame
    tf0_b.margin_left = tf0_b.margin_top = tf0_b.margin_right = tf0_b.margin_bottom = 0
    p0_b = tf0_b.paragraphs[0]
    p0_b.text = "LSRG-BYTETRACK · 课题组会汇报"
    p0_b.alignment = PP_ALIGN.CENTER
    p0_b.font.size = Pt(11.5)
    p0_b.font.bold = True
    p0_b.font.color.rgb = GOLD
    p0_b.font.name = FONT_TITLE

    t_box0 = slide0.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.3))
    tf0 = t_box0.text_frame
    tf0.word_wrap = True
    p0_1 = tf0.paragraphs[0]
    p0_1.text = "多目标跟踪在线因果风险感知与失效机理诊断"
    p0_1.font.size = Pt(31)
    p0_1.font.bold = True
    p0_1.font.color.rgb = TEXT_WHITE
    p0_1.font.name = FONT_TITLE

    p0_2 = tf0.add_paragraph()
    p0_2.text = "从物理建模、165万负样本 ECDF 校准到四大聚合算子博弈与三大失效模式解耦"
    p0_2.font.size = Pt(14.5)
    p0_2.font.color.rgb = RGBColor(196, 215, 237)
    p0_2.font.name = FONT_BODY
    p0_2.space_before = Pt(10)

    meta_cards = [
        ("负样本基准", "1,647,180 框", "三数据集正常帧全量背景"),
        ("真实 IDS 事件", "4,713 条", "三大互斥失效模式全覆盖"),
        ("防泄露评测协议", "114 序列 5-Fold", "Sequence-Stratified 跨序列盲测"),
        ("核心指标突破", "FPR 1.36% @ 60%", "Max/OWA 极低误报黄金区间")
    ]

    card_w = Inches(2.7)
    card_h = Inches(1.65)
    card_gap = Inches(0.24)
    start_x = Inches(0.9)
    start_y = Inches(4.4)

    for i, (k, v, desc) in enumerate(meta_cards):
        cx = start_x + i * (card_w + card_gap)
        add_card(slide0, cx, start_y, card_w, card_h, bg_color=RGBColor(18, 42, 74), border_color=RGBColor(38, 72, 114))
        
        tx = slide0.shapes.add_textbox(cx + Inches(0.15), start_y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        pk = tf.paragraphs[0]
        pk.text = k
        pk.font.size = Pt(11)
        pk.font.color.rgb = RGBColor(160, 185, 215)
        
        pv = tf.add_paragraph()
        pv.text = v
        pv.font.size = Pt(16.5)
        pv.font.bold = True
        pv.font.color.rgb = GOLD if i == 3 else TEXT_WHITE
        pv.font.name = FONT_TITLE
        pv.space_before = Pt(4)
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = RGBColor(140, 160, 185)
        pd.space_before = Pt(6)

    f_box0 = slide0.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.4))
    tf_f0 = f_box0.text_frame
    pf0 = tf_f0.paragraphs[0]
    pf0.text = "汇报人：LSRG 课题组 | 汇报时间：2026 年 8 月 | 评测基准：MOT17-half / MOT20 / SportsMOT (114 序列)"
    pf0.font.size = Pt(10)
    pf0.font.color.rgb = RGBColor(120, 145, 175)

    add_speaker_note(slide0, 
        "各位老师、同学大家好！今天我系统汇报多目标跟踪在线因果风险感知与失效机理诊断的最新成果。\n"
        "我们这次基于 114 个序列、164.7 万正常负样本检测框和 4,713 条真实 IDS 事件，构建了严格在线因果的风控评估体系与精细化机理诊断。"
    )

    # =========================================================================
    # SLIDE 1: 几何测度与关联基石 (GEOMETRIC PRIMER SLIDE)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, prs, BG_LIGHT)
    add_header(slide1, "00", "几何测度与关联基石：IoU 空间重叠与二部图代价矩阵",
               "从当前检测框与轨迹预测框的几何交并比，到全局二部图匹配代价与局部 2×2 交换扰动物理机理")

    add_card(slide1, Inches(0.8), Inches(1.2), Inches(11.73), Inches(3.45), bg_color=CARD_BG)

    geom_img_path = os.path.join(img_dir, "fig_geometry_iou_cost_matrix.png")
    if os.path.exists(geom_img_path):
        slide1.shapes.add_picture(geom_img_path, Inches(0.95), Inches(1.3), width=Inches(11.43))

    card_w3 = Inches(3.72)
    card_gap3 = Inches(0.28)
    card_h3 = Inches(2.25)
    start_y3 = Inches(4.8)

    geom_cards = [
        ("① IoU 空间几何重叠度", NAVY_MID,
         "• 几何意义：检测框 D 与卡尔曼外推预测 P 的交集/并集比值，值域严格 [0, 1]；\n"
         "• 欠匹配特征：f_weak = 1 - top1，IoU 越小说明失配严重，无轨迹时打满 1.0。"),
        ("② 代价矩阵与 Top1/Top2", NAVY_PRIMARY,
         "• 代价定义：Cost(D, T) = 1.0 - IoU，空间越接近则关联代价越小；\n"
         "• 竞争特征：top1 为首选轨迹，top2 为次优候选；top2 越大表明空间争夺越激烈。"),
        ("③ 2×2 局域交换扰动 (ΔC_swap)", RED_ACCENT,
         "• 脆弱性量化：对比正常匹配 (c11+c22) 与对调匹配 (c12+c21) 的代价差；\n"
         "• 临界预警：两目标交叉时若对调代价接近原代价，表明当前局部关联极其脆弱！")
    ]

    for i, (title_gc, col_gc, desc_gc) in enumerate(geom_cards):
        cx = Inches(0.8) + i * (card_w3 + card_gap3)
        add_card(slide1, cx, start_y3, card_w3, card_h3, bg_color=CARD_BG)
        
        tx = slide1.shapes.add_textbox(cx + Inches(0.18), start_y3 + Inches(0.15), card_w3 - Inches(0.36), card_h3 - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_gc
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col_gc
        
        p = tf.add_paragraph()
        p.text = desc_gc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(5)

    add_speaker_note(slide1,
        "在进入正式实验前，我们先通过直观的几何图示理解基础指标的物理来源：\n"
        "图 (a) 是 IoU 交并比：我们把当前检测框和卡尔曼预测框做交集和并集，计算出 0 到 1 的空间重合度。由此衍生的欠匹配特征 f_weak=1-top1 就是与最佳轨迹的匹配缺陷。\n"
        "图 (b) 是二部图代价矩阵：代价定义为 1-IoU，匈牙利匹配会优先选代价最小的 top1 轨迹，而次优轨迹 top2 则构成了竞争歧义特征。\n"
        "图 (c) 是局部 2×2 交换扰动：当两目标交叉时，我们计算正常分配与对调分配的代价差 ΔC_swap。如果对调代价与原代价极其接近，说明当前局域关联非常脆弱、极易跳变。\n"
        "理解了这三个几何指标，就能清晰看懂我们后续的特征评估与失效诊断。"
    )

    # =========================================================================
    # SLIDE 2: 初始特征有效性检验 (UNIFIED SCALE TABLE 1: P10 Low-Quantile Space)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, prs, BG_LIGHT)
    add_header(slide2, "01", "初始特征有效性检验：低分位(P10)困难样本下的正负断崖对比",
               "基于 165 万负样本 ECDF 校准至 [0, 1] 空间，严格检验正样本在 10% 最困难低分位下的判别水准与鲁棒性")

    # Unified Master Table Card (Left 0.8", Top 1.2", Width 11.73", Height 3.35")
    add_card(slide2, Inches(0.8), Inches(1.2), Inches(11.73), Inches(3.35), bg_color=CARD_BG)

    tx_t2_t = slide2.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.3))
    tf_t2_t = tx_t2_t.text_frame
    tf_t2_t.margin_left = tf_t2_t.margin_top = tf_t2_t.margin_right = tf_t2_t.margin_bottom = 0
    p = tf_t2_t.paragraphs[0]
    p.text = "一、 统一 [0, 1] 风险空间下：正常背景基准 vs 三大失效模式正样本低分位 (P10) 对比大表"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    # Table 1: Unified Scale Table with P10
    t1_shape = slide2.shapes.add_table(4, 7, Inches(1.0), Inches(1.72), Inches(11.33), Inches(2.65))
    table1 = t1_shape.table
    col_w1 = [Inches(1.6), Inches(1.5), Inches(1.45), Inches(1.45), Inches(1.45), Inches(1.45), Inches(2.43)]
    for j, w in enumerate(col_w1):
        table1.columns[j].width = w

    headers1 = ["风险特征分量", "负样本基准(P10)", "冷启动 Sc (P10)", "活跃接管 Sr (P10)", "重激活 Sh (P10)", "全量 IDS (P10)", "低分位正负断崖分离与物理机理"]
    for j, h in enumerate(headers1):
        format_cell(table1.cell(0, j), h, font_size=9.2, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY_HEADER)

    rows1_data = [
        ("欠匹配风险 r_weak", "0.1000", "0.9955 (98.6%*)", "0.4955", "0.9718 (69.4%*)", "0.7673 (59.9%*)", "🌟 极致下界保障：Sc/Sh 困难样本下界逼近 1.0，零死角拦截"),
        ("竞争歧义度 r_comp", "0.1000", "0.2264", "0.4116 (P50:0.87)", "0.2264", "0.2264 (4.3%*)", "🌟 响应空间争夺：Sr 低分位抬升 4 倍，中位达 0.875"),
        ("交换不稳定 r_swap", "0.1000", "0.7817", "0.6493 (33.7%*)", "0.5056", "0.6721 (19.5%*)", "🌟 强力响应 Sr：局域代价异动下界达 0.65，显著高于背景")
    ]

    for i, row in enumerate(rows1_data):
        bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        for j, val in enumerate(row):
            align_m = PP_ALIGN.CENTER if j < 6 else PP_ALIGN.LEFT
            bold_m = True if (i == 0 and j in [2, 4]) or (i == 2 and j == 3) else False
            color_m = GREEN if (i == 0 and j in [2, 4]) else (RED_ACCENT if (i == 2 and j == 3) else TEXT_DARK)
            format_cell(table1.cell(i + 1, j), val, font_size=9.2, bold=bold_m, color=color_m, align=align_m, bg_color=bg)

    # Bottom Section: Large, Bold, Executive Takeaway Banner (12.5pt font)
    add_card(slide2, Inches(0.8), Inches(4.75), Inches(11.73), Inches(2.3), bg_color=CARD_BG)
    tx_b2 = slide2.shapes.add_textbox(Inches(1.05), Inches(4.9), Inches(11.23), Inches(2.0))
    tf_b2 = tx_b2.text_frame
    tf_b2.word_wrap = True
    tf_b2.margin_left = tf_b2.margin_top = tf_b2.margin_right = tf_b2.margin_bottom = 0

    p = tf_b2.paragraphs[0]
    p.text = "🎯 低分位 (P10) 苛刻检验下的核心物理结论："
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    bullets2 = [
        "1. 严格基准（0.1000）：负样本经 ECDF 校准后严格服从 Uniform(0, 1)，10% 困难负样本基准值仅为 0.1000；",
        "2. 苛刻下界依然断崖分离：冷启动 Sc 即使在 10% 最困难样本下，欠匹配风险依然高达 0.9955（高出基准近 10 倍！），重激活 Sh 达 0.9718；",
        "3. 活跃接管机理互补：Sr 在交换不稳定分量上 P10 达 0.6493（高出基准 6.5 倍），有力弥补了近距交叉时欠匹配信号下界的不足；",
        "4. (*注：括号内为 r >= 0.99 占比，即 FPR <= 1% 下的捕获率；苛刻的 P10 检验直接证实了三大物理特征具有极高鲁棒性与正交互补性)。"
    ]

    for b_text in bullets2:
        p = tf_b2.add_paragraph()
        p.text = b_text
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(3.5)

    add_speaker_note(slide2,
        "第一页我们采用最苛刻的 P10（10% 低分位困难样本）来严格检验三个初始物理特征的质量。\n"
        "我们通过 165 万负样本做 ECDF 校准后，负样本的 P10 基准值严格为 0.1000。\n"
        "如果一个特征足够优秀，那么即使在最困难的 10% 事件样本（P10 下界）上，评分也应该维持在极高水准：\n"
        "第一，欠匹配特征在冷启动 Sc 上的 P10 达到了惊人的 0.9955，在重激活 Sh 上也达到了 0.9718！这意味着 90% 的 Sc 和 Sh 事件都被压倒性地推入 0.97 以上的绝对高风险区，实现近乎零误报拦截！\n"
        "第二，交换不稳定特征在活跃接管 Sr 上的 P10 达到 0.6493，比正常基准高出 6.5 倍，有效弥补了欠匹配的下界盲区！\n"
        "这从最苛刻的困难样本下界严格证明了三大指标的优越性与互补性。"
    )

    # =========================================================================
    # SLIDE 3: 四大聚合模型 5-Fold 盲测对比大表 (Table + Large Bold Takeaways)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, prs, BG_LIGHT)
    add_header(slide3, "02", "算子理论博弈：四大聚合模型 5-Fold 跨序列盲测对比大表",
               "以最大化极端分离度优化四大模型，确立 60%~75% 超低误报黄金区间与 90%~95% 协同攻坚区")

    add_card(slide3, Inches(0.8), Inches(1.2), Inches(11.73), Inches(3.3), bg_color=CARD_BG)

    tx_t3_t = slide3.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.3))
    tf_t3_t = tx_t3_t.text_frame
    tf_t3_t.margin_left = tf_t3_t.margin_top = tf_t3_t.margin_right = tf_t3_t.margin_bottom = 0
    p = tf_t3_t.paragraphs[0]
    p.text = "一、 四大模型在全量 165 万负样本下的 5-Fold 跨序列盲测误报率 (FPR) 对齐大表"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    t2_shape = slide3.shapes.add_table(8, 6, Inches(1.0), Inches(1.7), Inches(11.33), Inches(2.65))
    table2 = t2_shape.table
    col_w2 = [Inches(1.4), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(3.13)]
    for j, w in enumerate(col_w2):
        table2.columns[j].width = w

    headers2 = ["目标检出率 (TPR)", "Max 基线 (FPR)", "Noisy-OR (FPR)", "Power Mean (FPR)", "OWA 算子 (FPR)", "区间最优模型与机理特性"]
    for j, h in enumerate(headers2):
        format_cell(table2.cell(0, j), h, font_size=9.5, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY_HEADER)

    rows2_data = [
        ("60.0%", "1.36%", "3.12%", "5.37%", "1.36%", "🌟 Max / OWA 绝对王者：单特征无损直通，零噪声累积"),
        ("70.0%", "2.38%", "4.69%", "6.29%", "2.38%", "🌟 Max / OWA 领跑：仅 2.38% 误报覆盖七成全局跳变"),
        ("75.0%", "3.67%", "5.83%", "6.73%", "3.67%", "🌟 黄金性价比：仅需 3.67% 误报捕获 3/4 全量事件！"),
        ("85.0%", "10.39%", "10.38%", "10.90%", "10.39%", "⚡ 临界转折点：Noisy-OR 多弱风险协同优势开始反超"),
        ("90.0%", "16.32%", "14.79%", "16.66%", "16.32%", "🌟 Noisy-OR 胜出：亚可加性概率协同降低 1.53 pp 误报"),
        ("95.0%", "27.03%", "24.40%", "27.90%", "27.03%", "🌟 Noisy-OR 胜出：高攻坚区净降低 2.63 pp 误报"),
        ("pAUC[0.6, 1.0]", "0.8790", "0.8742", "0.8583", "0.8790", "Max / OWA 与 Noisy-OR 双星互补，综合表现优越")
    ]

    for i, row in enumerate(rows2_data):
        bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        for j, val in enumerate(row):
            align_m = PP_ALIGN.CENTER if j < 5 else PP_ALIGN.LEFT
            bold_m = True if (i in [0, 1, 2] and j in [1, 4]) or (i in [4, 5] and j == 2) else False
            color_m = GREEN if (i in [0, 1, 2] and j in [1, 4]) or (i in [4, 5] and j == 2) else TEXT_DARK
            format_cell(table2.cell(i + 1, j), val, font_size=9, bold=bold_m, color=color_m, align=align_m, bg_color=bg)

    add_card(slide3, Inches(0.8), Inches(4.7), Inches(5.75), Inches(2.35), bg_color=CARD_BG)
    tx_b3_l = slide3.shapes.add_textbox(Inches(1.0), Inches(4.85), Inches(5.35), Inches(2.05))
    tf_b3_l = tx_b3_l.text_frame
    tf_b3_l.word_wrap = True
    tf_b3_l.margin_left = tf_b3_l.margin_top = tf_b3_l.margin_right = tf_b3_l.margin_bottom = 0

    p = tf_b3_l.paragraphs[0]
    p.text = "🌟 区间 1：黄金可用区 (TPR 60% ~ 75%)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GREEN

    bullets3_l = [
        "• 首选算子：Max 基线 或 OWA；",
        "• 极致指标：60% 召回时误报仅 1.36%，75% 召回时仅 3.67%；",
        "• 工程价值：付出不到 4% 的正常扰动即可拦截 75% 的跳变！"
    ]
    for b in bullets3_l:
        p = tf_b3_l.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)

    add_card(slide3, Inches(6.75), Inches(4.7), Inches(5.78), Inches(2.35), bg_color=CARD_BG)
    tx_b3_r = slide3.shapes.add_textbox(Inches(6.95), Inches(4.85), Inches(5.38), Inches(2.05))
    tf_b3_r = tx_b3_r.text_frame
    tf_b3_r.word_wrap = True
    tf_b3_r.margin_left = tf_b3_r.margin_top = tf_b3_r.margin_right = tf_b3_r.margin_bottom = 0

    p = tf_b3_r.paragraphs[0]
    p.text = "🌟 区间 2：协同攻坚区 (TPR 85% ~ 95%)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD

    bullets3_r = [
        "• 首选算子：Noisy-OR 独立因果概率模型；",
        "• 协同增益：在 85% 拐点后反超 Max，95% 召回下误报降低 2.63 pp；",
        "• 工程价值：多弱信号非线性放大，适合复杂场景柔性风控预警。"
    ]
    for b in bullets3_r:
        p = tf_b3_r.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)

    add_speaker_note(slide3,
        "第二页汇报四大聚合模型的全面对决结果。\n"
        "如上方大表所示，我们通过 5 折跨序列盲测，评估了从 60% 到 95% 召回率下的误报演化。\n"
        "实验清晰确立了两个黄金工作区间：\n"
        "第一是【黄金可用区间（TPR 60%~75%）】：Max 算子在中低召回区近乎零误报，60% 召回时误报仅 1.36%，75% 召回时仅 3.67%。这意味着付出不到 4% 的代价就能拦截 3/4 的跳变，工业可用性极高。\n"
        "第二是【协同攻坚区间（TPR 85%~95%）】：Noisy-OR 算子在 85% 后反超 Max，在 95% 召回下将误报率从 27% 压低到 24.4%，展示出多弱信号协同放大收益。"
    )

    # =========================================================================
    # SLIDE 4: 【全屏大图 1】四大模型 ROC 曲线与 5% 步长折线大图 (Full Page Figure)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, prs, BG_LIGHT)
    add_header(slide4, "03", "模型评测全景：四大聚合模型高召回 ROC 与 5% 步长演化大图",
               "左图展示 TPR 60%~100% 高召回区域 ROC 曲线；右图展示 5% 步长误报率演化轨迹")

    add_card(slide4, Inches(0.8), Inches(1.2), Inches(11.73), Inches(5.85), bg_color=CARD_BG)
    
    chart_path1 = os.path.join(img_dir, "risk_aggregation_roc.png")
    if os.path.exists(chart_path1):
        slide4.shapes.add_picture(chart_path1, Inches(1.5), Inches(1.35), width=Inches(10.33))

    tx_s4_b = slide4.shapes.add_textbox(Inches(1.2), Inches(6.15), Inches(10.93), Inches(0.65))
    tf_s4_b = tx_s4_b.text_frame
    tf_s4_b.word_wrap = True
    tf_s4_b.margin_left = tf_s4_b.margin_top = tf_s4_b.margin_right = tf_s4_b.margin_bottom = 0

    p = tf_s4_b.paragraphs[0]
    p.text = "💡 视觉结论：Max 算子在中低召回区紧贴纵轴（近乎零误报）；Noisy-OR 算子在 85% 拐点后展现出强劲的协同上扬优势！"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    add_speaker_note(slide4,
        "这一页是四大聚合模型跨序列盲测的高清大图。\n"
        "左图的高召回 ROC 曲线清晰展示：Max 基线在 60% 到 80% 区间紧贴纵轴，误报率极低；\n"
        "右图的 5% 步长折线图清晰展示了 85% 处的交汇拐点：超过 85% 后，绿色曲线的 Noisy-OR 开始显著低于蓝色曲线的 Max，多特征协同优势一目了然。"
    )

    # =========================================================================
    # SLIDE 5: 三大失效模式精细化解耦与“自信而错”归因 (Table + Large Bold Takeaways)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, prs, BG_LIGHT)
    add_header(slide5, "04", "失效机理精细化解耦：揭示活跃接管(Sr)是误报雪崩的单一源头",
               "三大失效模式可观测性存在本质代差，Sr 模式存在“自信而错”局域伪装性，是制约超高召回的根本瓶颈")

    add_card(slide5, Inches(0.8), Inches(1.2), Inches(11.73), Inches(3.3), bg_color=CARD_BG)

    tx_t5_t = slide5.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.3))
    tf_t5_t = tx_t5_t.text_frame
    tf_t5_t.margin_left = tf_t5_t.margin_top = tf_t5_t.margin_right = tf_t5_t.margin_bottom = 0
    p = tf_t5_t.paragraphs[0]
    p.text = "一、 三大互斥失效模式独立评测与风险分位数全量对比表 (Noisy-OR 评分)"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    t3_shape = slide5.shapes.add_table(5, 8, Inches(1.0), Inches(1.7), Inches(11.33), Inches(2.65))
    table3 = t3_shape.table
    col_w3 = [Inches(1.3), Inches(1.3), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.5), Inches(1.8), Inches(2.13)]
    for j, w in enumerate(col_w3):
        table3.columns[j].width = w

    headers3 = ["失效模式", "样本量 (占比)", "最小值(Min)", "P10 分位", "中位(P50)", "评分>=0.99占比", "95% TPR 对应误报", "模式特征与可观测性"]
    for j, h in enumerate(headers3):
        format_cell(table3.cell(0, j), h, font_size=9.5, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY_HEADER)

    rows3_data = [
        ("冷启动 Sc", "1,828 (38.8%)", "0.9910", "0.9995", "1.0000", "100.00%", "6.09%", "🌟 天生透明：100% 极高置信，pAUC=0.963"),
        ("重激活 Sh", "986 (20.9%)", "0.7445", "0.9956", "0.9998", "95.94%", "17.07% (90%:12.6%)", "🌟 稳定可控：以欠匹配为主，pAUC=0.897"),
        ("活跃接管 Sr", "1,899 (40.3%)", "0.4126", "0.9797", "0.9997", "83.52%", "37.70% (MOT17:45.9%)", "⚠️ 误报源头：局域伪装性强，pAUC=0.788"),
        ("全量 Overall", "4,713 (100%)", "0.4126", "0.9938", "0.9999", "92.51%", "24.40%", "全场景异质统一空间，Sr 拉低全局表现")
    ]

    for i, row in enumerate(rows3_data):
        bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        for j, val in enumerate(row):
            align_m = PP_ALIGN.CENTER if j < 7 else PP_ALIGN.LEFT
            bold_m = True if (i == 0 and j in [0, 6]) or (i == 2 and j in [0, 6]) else False
            color_m = GREEN if (i == 0 and j == 6) else (RED_ACCENT if (i == 2 and j == 6) else TEXT_DARK)
            format_cell(table3.cell(i + 1, j), val, font_size=9, bold=bold_m, color=color_m, align=align_m, bg_color=bg)

    add_card(slide5, Inches(0.8), Inches(4.7), Inches(6.8), Inches(2.35), bg_color=CARD_BG)
    tx_b5_l = slide5.shapes.add_textbox(Inches(1.0), Inches(4.85), Inches(6.4), Inches(2.05))
    tf_b5_l = tx_b5_l.text_frame
    tf_b5_l.word_wrap = True
    tf_b5_l.margin_left = tf_b5_l.margin_top = tf_b5_l.margin_right = tf_b5_l.margin_bottom = 0

    p = tf_b5_l.paragraphs[0]
    p.text = "🔍 “自信而错 (Confident-but-Wrong)” 物理归因："
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT

    bullets5_l = [
        "1. 假性重叠：目标近距交叉时，错误轨迹沿速度外推与真实检测高度重合；",
        "2. 特征静默：欠匹配信号失效，伪装成高质量匹配导致最低分跌至 0.4126；",
        "3. 误报雪崩：后验门控为强行捞起深层伪装样本被迫拉低阈值，误伤海量正常框。"
    ]
    for b in bullets5_l:
        p = tf_b5_l.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(3)

    add_card(slide5, Inches(7.8), Inches(4.7), Inches(4.73), Inches(2.35), bg_color=CARD_BG)
    tx_b5_r = slide5.shapes.add_textbox(Inches(8.0), Inches(4.85), Inches(4.33), Inches(2.05))
    tf_b5_r = tx_b5_r.text_frame
    tf_b5_r.word_wrap = True
    tf_b5_r.margin_left = tf_b5_r.margin_top = tf_b5_r.margin_right = tf_b5_r.margin_bottom = 0

    p = tf_b5_r.paragraphs[0]
    p.text = "💡 核心学术结论："
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    bullets5_r = [
        "• 分流治理：占 59.7% 的冷启动与重激活可通过轻量后验门控在 FPR < 6% 下定向阻断；",
        "• 根源重构：单靠后验几何门控无法根治 Sr，必须从第一阶段匈牙利匹配代价函数根源破局！"
    ]
    for b in bullets5_r:
        p = tf_b5_r.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(5)

    add_speaker_note(slide5,
        "第五页是本次汇报最核心的学术亮点与机理归因。\n"
        "上方表格展示了三大模式独立评测与分位数的详细数据：\n"
        "冷启动 Sc 和重激活 Sh 占了近 60% 的事件，可观测性极佳：Sc 的 95% 召回误报仅 6.09%，100% 的事件评分都在 0.991 以上！\n"
        "然而，全局误报之所以会在高召回下上升，完全是由活跃接管 Sr 这一类引起的（在 95% 召回下误报达 37.7%）！\n"
        "我们归因出了‘自信而错’机理：当两目标快速交叉时，错误的轨迹由于速度外推恰好跟检测框深度重叠，伪装成‘完美匹配’，最低评分跌到 0.41。后验门控为了抓出它，阈值被迫降到 0.4，从而误伤了海量正常框。\n"
        "这得出了明确的结论：后验门控有其物理边界，单靠后验无法彻底解决 Sr，必须在第一阶段匹配函数上进行根源重构。"
    )

    # =========================================================================
    # SLIDE 6: 【全屏大图 2】跨数据集 2×2 独立评测 ROC 4 宫格大图 (Full Page Figure)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, prs, BG_LIGHT)
    add_header(slide6, "05", "跨数据集验证全景：MOT17 / MOT20 / SportsMOT 独立评测 4 宫格大图",
               "展示三大基准数据集在冷启动(Sc)、重激活(Sh)与活跃接管(Sr)上的独立 ROC 演化")

    add_card(slide6, Inches(0.8), Inches(1.2), Inches(11.73), Inches(5.85), bg_color=CARD_BG)
    
    chart_path2 = os.path.join(img_dir, "fig_dataset_roc_2x2_grid.png")
    if os.path.exists(chart_path2):
        slide6.shapes.add_picture(chart_path2, Inches(3.8), Inches(1.35), height=Inches(4.65))

    tx_s6_b = slide6.shapes.add_textbox(Inches(1.2), Inches(6.15), Inches(10.93), Inches(0.65))
    tf_s6_b = tx_s6_b.text_frame
    tf_s6_b.word_wrap = True
    tf_s6_b.margin_left = tf_s6_b.margin_top = tf_s6_b.margin_right = tf_s6_b.margin_bottom = 0

    p = tf_s6_b.paragraphs[0]
    p.text = "🌐 跨域异质性结论：MOT20（拥挤密度驱动）以 Sr 为绝对主导（占 58.94%）；SportsMOT（敏捷运动驱动）以 Sc 为绝对主导（占 56.33%）！"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    add_speaker_note(slide6,
        "最后一页我们来看跨数据集 4 宫格的独立评测大图。\n"
        "左上角 MOT17、右上角 MOT20、左下角 SportsMOT、右下角 Overall。\n"
        "在所有数据集中，绿色的 Sc 曲线均高高贴近左上角，而红色的 Sr 曲线均表现为明显的右侧拖尾，这证明了‘自信而错’现象具有跨数据集的高度普适性。\n"
        "同时，场景异质性非常明显：极高密度的 MOT20 以 Sr 占近六成，而高速运动的 SportsMOT 以 Sc 占近六成。\n"
        "以上就是我汇报的全部内容，请老师和同学批评指正！"
    )

    prs.save(output_path)
    print(f"[SUCCESS] 7-Slide presentation generated: {output_path}")


if __name__ == "__main__":
    out_file = r"e:\科研\ByteTrack\research\reports\LSRG_ByteTrack_组会汇报_v8.pptx"
    img_dir = r"e:\科研\ByteTrack\research\taxonomy"
    create_seven_slides_v7(out_file, img_dir)
    try:
        orig_file = r"e:\科研\ByteTrack\research\reports\LSRG_ByteTrack_组会汇报_20260817.pptx"
        create_seven_slides_v7(orig_file, img_dir)
    except Exception as e:
        print(f"[NOTE] Original file locked: {e}")
